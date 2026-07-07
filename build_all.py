"""Single build script — screenshots + PPTX + PDF for NSBM Market."""
import asyncio, os, base64

ARTIFACTS = '/home/uraniumx/.gemini/antigravity-ide/brain/0d0cede0-440c-4829-97fd-adf46d3be806'
BASE      = 'http://127.0.0.1:8000'

# ── SCREENSHOTS ───────────────────────────────────────────────────────────────
async def screenshots():
    from playwright.async_api import async_playwright
    shots = {}
    async with async_playwright() as p:
        browser = await p.chromium.launch(args=['--no-sandbox','--disable-dev-shm-usage'])
        ctx  = await browser.new_context(viewport={'width':1440,'height':900})
        page = await ctx.new_page()

        for name, url in [('home', f'{BASE}/index.php'),
                           ('details', f'{BASE}/details.php?id=1')]:
            await page.goto(url, wait_until='networkidle', timeout=20000)
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f'{ARTIFACTS}/{name}.png', full_page=False)
            shots[name] = f'{ARTIFACTS}/{name}.png'
            print(f'  {name} OK ({os.path.getsize(shots[name])//1024}KB)')

        # Seed cart then screenshot cart + checkout
        await page.goto(f'{BASE}/seed_cart.php', wait_until='domcontentloaded', timeout=10000)
        await page.wait_for_timeout(500)
        for name, url in [('cart', f'{BASE}/cart.php'),
                           ('checkout', f'{BASE}/checkout.php')]:
            await page.goto(url, wait_until='networkidle', timeout=12000)
            await page.wait_for_timeout(2000)
            await page.screenshot(path=f'{ARTIFACTS}/{name}.png', full_page=False)
            shots[name] = f'{ARTIFACTS}/{name}.png'
            print(f'  {name} OK ({os.path.getsize(shots[name])//1024}KB) [{page.url.split("/")[-1]}]')

        # Admin
        await page.goto(f'{BASE}/admin/login.php', wait_until='domcontentloaded', timeout=12000)
        await page.wait_for_timeout(1500)
        await page.screenshot(path=f'{ARTIFACTS}/admin.png', full_page=False)
        shots['admin'] = f'{ARTIFACTS}/admin.png'
        print(f'  admin OK ({os.path.getsize(shots["admin"])//1024}KB)')

        await page.fill('input[name="username"]','admin')
        await page.fill('input[name="password"]','admin123')
        await page.click('button[type="submit"]')
        await page.wait_for_timeout(2500)

        for name, url in [('dashboard', f'{BASE}/admin/dashboard.php'),
                           ('orders',    f'{BASE}/admin/orders.php')]:
            await page.goto(url, wait_until='networkidle', timeout=15000)
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f'{ARTIFACTS}/{name}.png', full_page=False)
            shots[name] = f'{ARTIFACTS}/{name}.png'
            print(f'  {name} OK ({os.path.getsize(shots[name])//1024}KB)')

        await browser.close()
    return shots

# ── PPTX ─────────────────────────────────────────────────────────────────────
def build_pptx(shots):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from PIL import Image as PILImage

    GREEN  = RGBColor(0x00,0x68,0x37); LGREEN = RGBColor(0x00,0xa6,0x51)
    WHITE  = RGBColor(0xFF,0xFF,0xFF); DARK   = RGBColor(0x1a,0x1a,0x1a)
    LGRAY  = RGBColor(0xf0,0xf7,0xf0)

    def blank(prs):   return prs.slides.add_slide(prs.slide_layouts[6])
    def box(sl,l,t,w,h,fill=None):
        s=sl.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
        s.line.fill.background()
        if fill: s.fill.solid(); s.fill.fore_color.rgb=fill
        else:    s.fill.background()
        return s
    def TB(sl,l,t,w,h):
        b=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
        b.text_frame.word_wrap=True; return b.text_frame
    def R(tf,text,sz,bold=False,color=DARK,align=PP_ALIGN.LEFT):
        p=tf.paragraphs[0] if not tf.text else tf.add_paragraph()
        p.alignment=align; r=p.add_run()
        r.text=text; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=color
        return p
    def RR(tf,text,sz=13,bold=False,color=DARK,space=0):
        p=tf.add_paragraph(); p.space_before=Pt(space)
        r=p.add_run(); r.text=text; r.font.size=Pt(sz)
        r.font.bold=bold; r.font.color.rgb=color; return p
    def HDR(prs,title,speaker=''):
        sl=blank(prs); box(sl,0,0,13.33,1.5,GREEN); box(sl,0,1.5,13.33,0.06,LGREEN); box(sl,0,1.56,13.33,5.94,WHITE)
        f=TB(sl,0.4,0.18,10,1.15); R(f,title,28,bold=True,color=WHITE)
        if speaker: f2=TB(sl,9.5,0.22,3.5,0.5); R(f2,speaker,10,color=RGBColor(0xcc,0xff,0xcc),align=PP_ALIGN.RIGHT)
        return sl
    def TBLH(sl,cols,cx,cw,y=1.88,h=0.4):
        for x,w,t in zip(cx,cw,cols):
            box(sl,x,y,w,h,GREEN); f=TB(sl,x+0.08,y+0.06,w-0.12,h-0.08); R(f,t,11,bold=True,color=WHITE)
    def TBLR(sl,vals,cx,cw,y,h=0.52,even=True,last=False,fg=False):
        bg=GREEN if last else (LGRAY if even else WHITE)
        for i,(x,w,t) in enumerate(zip(cx,cw,vals)):
            box(sl,x,y,w,h,bg); f=TB(sl,x+0.08,y+0.07,w-0.12,h-0.1)
            clr=WHITE if last else (GREEN if(i==0 and fg) else DARK)
            R(f,t,11,bold=(last or(i==0 and fg)),color=clr)
    def IMG(sl,path,mw=12.9,mh=6.4,top=0.62):
        if not os.path.exists(path): return
        with PILImage.open(path) as im: iw,ih=im.size
        sc=min(mw/iw,mh/ih); dw,dh=iw*sc,ih*sc
        sl.shapes.add_picture(path,Inches((13.33-dw)/2),Inches(top+(mh-dh)/2),Inches(dw),Inches(dh))

    prs=Presentation(); prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)

    # COVER
    sl=blank(prs); box(sl,0,0,13.33,7.5,GREEN); box(sl,0,0,0.14,7.5,LGREEN)
    f=TB(sl,0.5,1.0,12,0.38); R(f,'NSBM GREEN UNIVERSITY · FACULTY OF COMPUTING',9,color=RGBColor(0xbb,0xff,0xbb))
    f=TB(sl,0.5,1.55,12,1.4); R(f,'NSBM Market',58,bold=True,color=WHITE)
    f=TB(sl,0.5,3.05,12,0.52); R(f,'Web Based Application Development — Final Project',17,color=RGBColor(0xcc,0xff,0xcc))
    box(sl,0.5,3.75,5,0.05,RGBColor(0x33,0xaa,0x66))
    f=TB(sl,0.5,3.95,12,2.7)
    R(f,'34473  KADK Gimhana Kumarapeli  (Team Leader)',13,bold=True,color=WHITE)
    for m in ['35550  KTSR Viduranga','35546  DKR Wijethunge','34467  SERI Senewirathna','34673  AK Nethsara']:
        RR(f,m,13,color=RGBColor(0xdd,0xff,0xdd))
    f2=TB(sl,10.5,7.05,2.5,0.35); R(f2,'July 2026',10,color=RGBColor(0x88,0xcc,0x88),align=PP_ALIGN.RIGHT)

    # PROJECT SUMMARY
    sl=HDR(prs,'Project Summary','Gimhana (34473)')
    f=TB(sl,0.4,1.72,12.5,5.5)
    R(f,'Campus e-commerce platform — NSBM students and staff',15,bold=True,color=GREEN)
    RR(f,'HTML · CSS · Vanilla JavaScript · PHP · MySQL  —  no heavy frameworks',13)
    RR(f,'',5); RR(f,'Features:',14,bold=True,color=GREEN)
    for feat in ['  ✓  AJAX product filtering — no page reloads',
                 '  ✓  Session-based cart with variant tracking',
                 '  ✓  Pre-order product system (10 products)',
                 '  ✓  Admin panel: full CRUD + multi-image upload',
                 '  ✓  Dark / Light mode (localStorage persistence)',
                 '  ✓  PDO prepared statements — SQL injection proof']:
        RR(f,feat,13)
    RR(f,'',8)
    RR(f,'📁  drive.google.com/drive/folders/1gtQQGC3yyVT1MIzuk4-Z5eEwiLfuP4gf',11,color=GREEN)
    RR(f,'📞  Team Leader Contact: 0766060863',11,color=GREEN)

    # TECH STACK
    sl=HDR(prs,'Technology Stack','Gimhana (34473)')
    cx=[0.4,2.1,6.3]; cw=[1.6,4.1,6.6]
    TBLH(sl,['Layer','Technology','Rationale'],cx,cw)
    for ri,v in enumerate([('Frontend','HTML5, Vanilla CSS, Vanilla JavaScript','Module requirement — no frameworks'),
        ('Backend','PHP (Procedural)','Mandatory per module specification'),
        ('Database','MySQL / MariaDB via XAMPP','Mandatory per module specification'),
        ('DB Access','PDO + prepared statements','SQL-injection proof by design'),
        ('Hosting','XAMPP local / InfinityFree production','Auto-detected via socket in config/db.php')]):
        TBLR(sl,v,cx,cw,y=2.28+ri*0.56,even=(ri%2==0))

    # ARCHITECTURE
    sl=HDR(prs,'System Architecture','Gimhana (34473)')
    f=TB(sl,0.4,1.72,12.5,5.6)
    R(f,'Customer Layer',14,bold=True,color=GREEN)
    for x in ['  index.php — Storefront + AJAX product grid + sliding filter sidebar',
               '  fetch_products.php — AJAX endpoint: dynamic PDO query → returns HTML cards',
               '  details.php — Image gallery, variant chips, Add to Cart / Pre-Order Now',
               '  cart.php + cart_action.php — Session cart (product_id–variant_id key)',
               '  checkout.php — Order form → inserts into orders + order_items']: RR(f,x,12)
    RR(f,'',6); RR(f,'Admin Layer',14,bold=True,color=GREEN)
    for x in ['  login.php → session guard on every admin page',
               '  dashboard.php → product listing + edit / delete',
               '  add_product.php / edit_product.php → multi-image upload (first = thumbnail)',
               '  orders.php → view all orders with status']: RR(f,x,12)
    RR(f,'',6); RR(f,'Shared: config/db.php (auto local/prod)  ·  image_helper.php (local path resolver)',12,bold=True,color=GREEN)

    # CATALOGUE
    sl=HDR(prs,'Product Catalogue  —  80 Products','Ravindu (35546)')
    cx=[0.4,2.75,3.65]; cw=[2.25,0.8,9.0]
    TBLH(sl,['Category','Count','Sample Items'],cx,cw)
    for ri,v in enumerate([('Apparel','16','NSBM Hoodie, Denim Jacket, Windbreaker, Cargo Pants'),
        ('Stationery','16','Customized Notebook, Fountain Pen, Acrylic Paint Set'),
        ('Tech & Gadgets','16','Bluetooth Earbuds, External SSD, Smart Watch, Laptop Stand'),
        ('Accessories','16','Yoga Mat, Travel Mug, Bento Lunch Box, Phone Case'),
        ('Food & Snacks','16','Homemade Brownies, Dark Chocolate Box, Matcha Kit'),
        ('TOTAL','80','10 Pre-Order  ·  70 Active / In Stock')]):
        TBLR(sl,v,cx,cw,y=2.28+ri*0.56,even=(ri%2==0),last=(ri==5))

    # KEY FEATURES
    sl=HDR(prs,'Key Technical Features','Gimhana (34473)')
    for i,(title,desc) in enumerate([
        ('⚡  AJAX Filtering','fetch_products.php builds a dynamic parameterised PDO query from filter params and returns HTML cards. Zero page reloads.'),
        ('🛒  Session Cart','$_SESSION["cart"] uses product_id–variant_id composite key. Same product in 2 sizes tracked independently. Nav badge via AJAX.'),
        ('🕐  Pre-Order System','status ENUM(active,pre_order,out_of_stock) in DB. Orange badge on catalogue cards. PRE-ORDER NOW button on details page.'),
        ('🌙  Dark / Light Mode','Moon/sun button in navbar toggles data-theme on <html>. CSS custom properties flip instantly. Preference saved to localStorage.')]):
        y=1.76+i*1.34; box(sl,0.4,y+0.05,0.07,1.05,GREEN)
        f=TB(sl,0.62,y,12.3,0.46); R(f,title,15,bold=True,color=GREEN)
        f2=TB(sl,0.62,y+0.46,12.3,0.76); R(f2,desc,12)

    # DB SCHEMA
    sl=HDR(prs,'Database Schema','Ravindu (35546)')
    cx=[0.4,3.25]; cw=[2.75,9.75]
    TBLH(sl,['Table','Key Columns'],cx,cw)
    for ri,v in enumerate([('products','id, name, price, category_id, image_path, status, stock_quantity'),
        ('product_variants','product_id, variant_name, variant_type, price_modifier, stock_quantity'),
        ('categories','id, name'),
        ('orders','id, customer_name, customer_email, total_amount, status'),
        ('order_items','order_id, product_id, variant_id, quantity, price'),
        ('users','id, username, password')]):
        TBLR(sl,v,cx,cw,y=2.28+ri*0.56,even=(ri%2==0),fg=True)

    # SETUP
    sl=HDR(prs,'Setup Instructions','Nethsara (34673)')
    f=TB(sl,0.4,1.72,12.5,5.5); R(f,'XAMPP Setup',15,bold=True,color=GREEN)
    for s in ['1.  Copy store/ into XAMPP htdocs/',
              '2.  Start Apache + MySQL from XAMPP Control Panel',
              '3.  phpMyAdmin → create database  nsbm_market',
              '4.  Import  sql/nsbm_market.sql  (creates tables + 80 products + 20 orders)',
              '5.  Open  http://localhost/store/',
              '6.  Admin: /admin/login.php  |  user: admin  |  password: admin123']: RR(f,s,13)
    RR(f,'',10); RR(f,'Project Structure',15,bold=True,color=GREEN)
    RR(f,'admin/ (9 files)  ·  assets/css/style.css  ·  assets/js/script.js',12)
    RR(f,'assets/images/products/ (80 local JPGs)  ·  config/db.php  ·  sql/nsbm_market.sql',12)
    RR(f,'index.php  ·  details.php  ·  cart.php  ·  checkout.php  ·  fetch_products.php  ·  cart_action.php',12)

    # SCREENSHOT SLIDES
    for key,title,speaker,note in [
        ('home','Home — Product Catalogue','Gimhana (34473)','AJAX grid · Filter sidebar: category, price, availability · Green=In Stock, Orange=Pre-Order'),
        ('details','Product Details Page','Viduranga (35550)','Thumbnail gallery · Variant chips · Stock status · Add to Cart / Pre-Order Now'),
        ('cart','Shopping Cart','Senewirathna (34467)','Session cart · Variant display · Quantity controls · Live AJAX badge update'),
        ('checkout','Checkout Flow','Senewirathna (34467)','Order form: name, email, phone, address · Inserts order + clears cart on submit'),
        ('dashboard','Admin Dashboard','Ravindu (35546)','80 products listed · Edit / Delete per row · Add Product for new listings'),
        ('orders','Admin Orders','Ravindu (35546)','20 seed orders + live orders · Customer name, email, total, status, date'),
        ('admin','Admin Login','Nethsara (34673)','Session-based auth · Admin URLs redirect here without valid session')]:
        sl=blank(prs); box(sl,0,0,13.33,0.57,GREEN); box(sl,0,7.1,13.33,0.4,LGRAY)
        f=TB(sl,0.3,0.07,9.5,0.44); R(f,title,18,bold=True,color=WHITE)
        f2=TB(sl,9.9,0.1,3.1,0.38); R(f2,speaker,10,color=RGBColor(0xcc,0xff,0xcc),align=PP_ALIGN.RIGHT)
        IMG(sl, shots.get(key,f'{ARTIFACTS}/{key}.png'))
        f3=TB(sl,0.3,7.13,12.8,0.35); R(f3,note,10,color=RGBColor(0x44,0x44,0x44))

    # CLOSING
    sl=blank(prs); box(sl,0,0,13.33,7.5,GREEN); box(sl,0,0,0.14,7.5,LGREEN)
    f=TB(sl,0.5,1.8,12,1.3); R(f,'Thank You',56,bold=True,color=WHITE)
    f2=TB(sl,0.5,3.2,12,0.5); R(f2,'NSBM Market — Web Based Application Development',16,color=RGBColor(0xcc,0xff,0xcc))
    f3=TB(sl,0.5,4.15,12,2.5)
    R(f3,'📁  drive.google.com/drive/folders/1gtQQGC3yyVT1MIzuk4-Z5eEwiLfuP4gf',13,color=WHITE)
    RR(f3,'📞  Team Leader Contact: 0766060863',13,color=WHITE)
    RR(f3,'',8); RR(f3,'HTML · CSS · Vanilla JS · PHP · MySQL · PDO · AJAX',12,color=RGBColor(0xaa,0xee,0xaa))

    out='/opt/lampp/htdocs/store/NSBM_Market_Presentation.pptx'
    prs.save(out)
    print(f'PPTX → {out}  ({os.path.getsize(out)//1024}KB)')

# ── PDF ───────────────────────────────────────────────────────────────────────
async def build_pdf(shots):
    import asyncio as _asyncio, base64 as _b64
    def b64(p):
        return _b64.b64encode(open(p,'rb').read()).decode() if p and os.path.exists(p) else None

    ss=[('Home — Product Catalogue','home.png'),('Product Details Page','details.png'),
        ('Shopping Cart','cart.png'),('Checkout Flow','checkout.png'),
        ('Admin Dashboard','dashboard.png'),('Admin Orders','orders.png'),('Admin Login','admin.png')]
    img_blocks=''.join(
        f'<div class="ss-wrap"><div class="ss-label">{cap}</div><img src="data:image/png;base64,{b64(f"{ARTIFACTS}/{fn}")}" alt="{cap}"></div>'
        for cap,fn in ss if b64(f'{ARTIFACTS}/{fn}')
    )
    DRIVE='https://drive.google.com/drive/folders/1gtQQGC3yyVT1MIzuk4-Z5eEwiLfuP4gf?usp=sharing'
    html=f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>
  @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;900&display=swap');
  *{{box-sizing:border-box;margin:0;padding:0}}
  body{{font-family:'Inter',Arial,sans-serif;color:#1a1a1a;font-size:12px;line-height:1.65;background:white}}
  .cover{{background:linear-gradient(135deg,#006837 0%,#00a651 100%);color:white;padding:70px 60px 60px;min-height:280px;display:flex;flex-direction:column;justify-content:flex-end}}
  .cover-tag{{font-size:.75rem;font-weight:600;letter-spacing:3px;text-transform:uppercase;opacity:.75;margin-bottom:14px}}
  .cover-title{{font-size:3.2rem;font-weight:900;letter-spacing:-2px;line-height:1;margin-bottom:12px}}
  .cover-sub{{font-size:1rem;font-weight:500;opacity:.85;margin-bottom:30px}}
  .cover-team{{display:flex;flex-direction:column;gap:4px;font-size:.8rem;opacity:.9;border-top:1px solid rgba(255,255,255,.3);padding-top:20px}}
  .content{{padding:50px 60px}}
  h1{{font-size:1.3rem;font-weight:900;color:#006837;text-transform:uppercase;letter-spacing:.5px;border-bottom:3px solid #006837;padding-bottom:8px;margin:45px 0 18px}}
  h1:first-child{{margin-top:0}}
  h2{{font-size:1rem;font-weight:700;color:#222;margin:28px 0 10px}}
  h3{{font-size:.9rem;font-weight:700;color:#006837;margin:18px 0 6px}}
  p{{margin-bottom:10px;color:#333}}
  table{{width:100%;border-collapse:collapse;margin:16px 0;font-size:.9rem}}
  thead th{{background:#006837;color:white;padding:9px 12px;text-align:left;font-weight:700;font-size:.8rem;text-transform:uppercase}}
  tbody td{{padding:8px 12px;border-bottom:1px solid #eef0ee;vertical-align:top}}
  tbody tr:nth-child(even) td{{background:#f8faf8}}
  pre{{background:#f4f7f4;border-left:3px solid #006837;padding:18px 20px;border-radius:6px;font-size:.78rem;line-height:1.6;white-space:pre-wrap;margin:16px 0;font-family:monospace}}
  code{{background:#edf2ed;padding:2px 5px;border-radius:3px;font-size:.85em;font-family:monospace}}
  ul,ol{{padding-left:22px;margin-bottom:12px}} li{{margin-bottom:4px;color:#333}}
  .info-box{{background:#f0f7f0;border:1px solid #c3e0c3;border-radius:8px;padding:18px 22px;margin:18px 0;font-size:.9rem;line-height:2}}
  .info-box .label{{font-weight:700;color:#006837;display:inline-block;min-width:140px}}
  .info-box a{{color:#006837;word-break:break-all}}
  .ss-wrap{{margin:28px 0;page-break-inside:avoid}}
  .ss-label{{font-weight:700;font-size:.8rem;color:#006837;text-transform:uppercase;letter-spacing:.8px;margin-bottom:8px;display:flex;align-items:center;gap:8px}}
  .ss-label::before{{content:'';display:inline-block;width:4px;height:16px;background:#006837;border-radius:2px}}
  .ss-wrap img{{width:100%;border-radius:8px;box-shadow:0 3px 14px rgba(0,0,0,.12);border:1px solid #e0e0e0;display:block}}
  .page-break{{page-break-before:always;height:0}}
  @page{{size:A4;margin:0}}
</style></head><body>
<div class="cover">
  <div class="cover-tag">NSBM Green University · Faculty of Computing</div>
  <div class="cover-title">NSBM Market</div>
  <div class="cover-sub">Web Based Application Development — Final Project Report</div>
  <div class="cover-team">
    <span>34473 &nbsp; KADK Gimhana Kumarapeli &nbsp; (Team Leader)</span>
    <span>35550 &nbsp; KTSR Viduranga</span>
    <span>35546 &nbsp; DKR Wijethunge</span>
    <span>34467 &nbsp; SERI Senewirathna</span>
    <span>34673 &nbsp; AK Nethsara</span>
    <div style="margin-top:12px;opacity:.7">July 2026</div>
  </div>
</div>
<div class="content">
<h1>1. Project Summary</h1>
<p>NSBM Market is a full-stack e-commerce web application built for NSBM Green University, enabling student entrepreneurs to sell campus merchandise — hoodies, stationery, tech accessories and food items — directly to NSBM students and staff. The platform features a dynamic AJAX product catalogue, session-based cart, pre-order system and a fully functional admin panel.</p>
<div class="info-box">
  <div><span class="label">Team Leader:</span> KADK Gimhana Kumarapeli (34473)</div>
  <div><span class="label">Contact:</span> 0766060863</div>
  <div><span class="label">Google Drive:</span> <a href="{DRIVE}">{DRIVE}</a></div>
  <div style="font-size:.8rem;color:#555;margin-top:6px">The Drive folder contains the full source code and video demonstration.</div>
</div>
<h1>2. Technology Stack</h1>
<table><thead><tr><th>Layer</th><th>Technology</th><th>Rationale</th></tr></thead><tbody>
<tr><td>Frontend</td><td>HTML5, Vanilla CSS, Vanilla JavaScript</td><td>Module requirement — no frameworks</td></tr>
<tr><td>Backend</td><td>PHP (Procedural)</td><td>Mandatory per module</td></tr>
<tr><td>Database</td><td>MySQL / MariaDB via XAMPP</td><td>Mandatory per module</td></tr>
<tr><td>DB Access</td><td>PDO with prepared statements</td><td>SQL-injection proof</td></tr>
<tr><td>Hosting</td><td>XAMPP (local) / InfinityFree (production)</td><td>Auto-detected in config/db.php</td></tr>
</tbody></table>
<h1>3. System Architecture</h1>
<h2>Customer-Facing Pages</h2>
<ul>
<li><code>index.php</code> — Storefront with AJAX product grid and sliding filter sidebar</li>
<li><code>fetch_products.php</code> — AJAX endpoint: builds dynamic PDO query, returns HTML cards</li>
<li><code>details.php</code> — Product page with gallery, variant chips, cart/pre-order actions</li>
<li><code>cart.php + cart_action.php</code> — Session cart (product_id–variant_id composite key)</li>
<li><code>checkout.php</code> — Order placement, inserts to orders and order_items</li>
</ul>
<h2>Admin Panel</h2>
<ul>
<li><code>admin/login.php</code> — Session-based authentication, guards all admin pages</li>
<li><code>admin/dashboard.php</code> — Product list with edit/delete</li>
<li><code>admin/add_product.php + edit_product.php</code> — Multi-image upload (first = thumbnail)</li>
<li><code>admin/orders.php</code> — All customer orders with status tracking</li>
</ul>
<h1>4. Product Catalogue</h1>
<table><thead><tr><th>Category</th><th>Products</th><th>Sample Items</th></tr></thead><tbody>
<tr><td>Apparel</td><td>16</td><td>NSBM Hoodie, Denim Jacket, Windbreaker, Cargo Pants</td></tr>
<tr><td>Stationery</td><td>16</td><td>Customized Notebook, Fountain Pen, Acrylic Paint Set</td></tr>
<tr><td>Tech &amp; Gadgets</td><td>16</td><td>Bluetooth Earbuds, External SSD, Smart Watch, Laptop Stand</td></tr>
<tr><td>Accessories</td><td>16</td><td>Yoga Mat, Travel Mug, Bento Lunch Box, Phone Case</td></tr>
<tr><td>Food &amp; Snacks</td><td>16</td><td>Homemade Brownies, Dark Chocolate Box, Matcha Kit</td></tr>
<tr><td><strong>Total</strong></td><td><strong>80</strong></td><td>10 Pre-Order · 70 Active/In Stock</td></tr>
</tbody></table>
<h1>5. Key Technical Features</h1>
<h3>AJAX-Powered Filtering</h3>
<p>Products load and filter dynamically using the Fetch API. Users filter by category, price range, search keyword and availability. <code>fetch_products.php</code> builds a parameterised PDO query and returns rendered HTML cards — no page reload.</p>
<h3>Session-Based Cart</h3>
<p>Cart items stored in <code>$_SESSION['cart']</code> with a composite <code>product_id-variant_id</code> key so the same product in two sizes coexists independently. The nav cart badge updates in real-time via AJAX.</p>
<h3>Pre-Order System</h3>
<p>Products with <code>status = 'pre_order'</code> show an orange PRE-ORDER badge on catalogue cards and a PRE-ORDER NOW button on the details page instead of Add to Cart.</p>
<h3>Dark / Light Mode</h3>
<p>A moon/sun button in the navbar toggles the <code>data-theme</code> attribute on the html element, swapping all CSS custom properties instantly. Preference persists in <code>localStorage</code>.</p>
<h1>6. Database Schema</h1>
<table><thead><tr><th>Table</th><th>Key Columns</th><th>Purpose</th></tr></thead><tbody>
<tr><td><code>products</code></td><td>id, name, price, category_id, image_path, status, stock_quantity</td><td>Core product listings</td></tr>
<tr><td><code>product_variants</code></td><td>product_id, variant_name, variant_type, price_modifier, stock_quantity</td><td>Size, colour, capacity variants</td></tr>
<tr><td><code>categories</code></td><td>id, name</td><td>Five product categories</td></tr>
<tr><td><code>orders</code></td><td>id, customer_name, customer_email, total_amount, status</td><td>Customer orders</td></tr>
<tr><td><code>order_items</code></td><td>order_id, product_id, variant_id, quantity, price</td><td>Line items per order</td></tr>
<tr><td><code>users</code></td><td>id, username, password</td><td>Admin accounts</td></tr>
</tbody></table>
<h1>7. Setup Instructions</h1>
<ol>
<li>Copy <code>store/</code> folder into XAMPP <code>htdocs/</code>.</li>
<li>Start Apache and MySQL from XAMPP Control Panel.</li>
<li>phpMyAdmin → create database <code>nsbm_market</code> → import <code>sql/nsbm_market.sql</code>.</li>
<li>Open <code>http://localhost/store/</code>.</li>
<li>Admin panel: <code>/admin/login.php</code> — username: <code>admin</code> / password: <code>admin123</code>.</li>
</ol>
<h1>8. Project File Structure</h1>
<pre>store/
├── admin/           login, signup, logout, dashboard, add/edit/delete_product, orders, forgot_password
├── assets/
│   ├── css/style.css
│   ├── js/script.js
│   └── images/products/   (80 local product images)
├── config/db.php    PDO connection — auto local/production detection
├── includes/        header.php, footer.php, image_helper.php
├── sql/nsbm_market.sql    Schema + 80 products + 20 seed orders
├── index.php        Storefront
├── details.php      Product detail
├── cart.php         Cart
├── checkout.php     Order placement
├── fetch_products.php     AJAX product loader
└── cart_action.php  Cart handler</pre>
<div class="page-break"></div>
<h1>9. Screenshots</h1>
{img_blocks}
</div></body></html>'''

    tmp='/tmp/nsbm_report.html'
    with open(tmp,'w',encoding='utf-8') as f: f.write(html)

    async def pdf():
        from playwright.async_api import async_playwright
        async with async_playwright() as p:
            browser=await p.chromium.launch(args=['--no-sandbox'])
            page=await browser.new_page()
            await page.goto(f'file://{tmp}',wait_until='domcontentloaded')
            await page.wait_for_timeout(1500)
            out='/opt/lampp/htdocs/store/NSBM_Market_Project_Report.pdf'
            await page.pdf(path=out,format='A4',print_background=True,margin={'top':'0px','right':'0px','bottom':'0px','left':'0px'})
            await browser.close()
        os.remove(tmp)
        print(f'PDF → {out}  ({os.path.getsize(out)//1024}KB)')

    await pdf()

async def main():
    # Create seed_cart.php temporarily
    seed = '''\
<?php
session_start();
$_SESSION['cart'] = [
    '1-0'  => ['product_id'=>1, 'variant_id'=>0,'quantity'=>2,'name'=>'NSBM Hoodie - Green','price'=>3500.00,'image'=>'assets/images/products/hoodie.jpg'],
    '34-0' => ['product_id'=>34,'variant_id'=>0,'quantity'=>1,'name'=>'Bluetooth Earbuds','price'=>5500.00,'image'=>'assets/images/products/earbuds.jpg'],
    '17-0' => ['product_id'=>17,'variant_id'=>0,'quantity'=>1,'name'=>'Customized Notebook','price'=>450.00,'image'=>'assets/images/products/notebook.jpg'],
];
header('Location: checkout.php');
exit;
?>'''
    with open('/opt/lampp/htdocs/store/seed_cart.php','w') as f: f.write(seed)

    print('Taking screenshots...')
    shots = await screenshots()
    os.remove('/opt/lampp/htdocs/store/seed_cart.php')

    print('Building PPTX...')
    build_pptx(shots)

    print('Building PDF...')
    await build_pdf(shots)
    print('All done.')

asyncio.run(main())
